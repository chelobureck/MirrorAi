"""
Роутер для экспорта презентаций в различные форматы
"""
from fastapi import APIRouter, Depends, HTTPException, status
from fastapi.responses import FileResponse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from models.base import get_session
from models.user import User
from models.presentation import Presentation
from utils.auth import get_current_user
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import tempfile
import os
from typing import Optional

router = APIRouter(prefix="/export", tags=["export"])

@router.get("/presentations/{presentation_id}/pptx")
async def export_to_pptx(
    presentation_id: int,
    current_user: User = Depends(get_current_user),
    session: AsyncSession = Depends(get_session)
):
    """Экспорт презентации в PPTX формат"""
    
    # Получаем презентацию
    result = await session.execute(
        select(Presentation).where(
            Presentation.id == presentation_id,
            Presentation.user_id == current_user.id
        )
    )
    presentation = result.scalars().first()
    
    if not presentation:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Презентация не найдена"
        )
    
    # Создаем PPTX
    pptx_presentation = PPTXPresentation()
    
    # Удаляем пустой слайд по умолчанию
    slide_layouts = pptx_presentation.slide_layouts
    
    for slide_data in presentation.content.get("slides", []):
        # Выбираем layout в зависимости от типа слайда
        if slide_data.get("type") == "title":
            slide_layout = slide_layouts[0]  # Title slide
        else:
            slide_layout = slide_layouts[1]  # Title and content
        
        slide = pptx_presentation.slides.add_slide(slide_layout)
        
        # Добавляем заголовок
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")
        
        # Добавляем содержимое
        if len(slide.shapes.placeholders) > 1 and slide_data.get("content"):
            content_placeholder = slide.shapes.placeholders[1]
            content_placeholder.text = slide_data.get("content", "")
    
    # Сохраняем во временный файл
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    pptx_presentation.save(temp_file.name)
    temp_file.close()
    
    # Возвращаем файл
    safe_title = "".join(c for c in presentation.title if c.isalnum() or c in (' ', '-', '_')).rstrip()
    filename = f"{safe_title}.pptx"
    
    return FileResponse(
        temp_file.name,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        filename=filename,
        background=lambda: os.unlink(temp_file.name)
    )

@router.get("/presentations/{presentation_id}/pdf")
async def export_to_pdf(
    presentation_id: int,
    current_user: User = Depends(get_current_user),
    session: AsyncSession = Depends(get_session)
):
    """Экспорт презентации в PDF формат"""
    
    # Получаем презентацию
    result = await session.execute(
        select(Presentation).where(
            Presentation.id == presentation_id,
            Presentation.user_id == current_user.id
        )
    )
    presentation = result.scalars().first()
    
    if not presentation:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Презентация не найдена"
        )
    
    # Создаем PDF
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(temp_file.name, pagesize=A4)
    
    # Стили
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        alignment=1  # Центрирование
    )
    
    slide_title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading2'],
        fontSize=18,
        spaceAfter=20,
        textColor='blue'
    )
    
    content_style = ParagraphStyle(
        'Content',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=15,
        leftIndent=20
    )
    
    story = []
    
    # Титульная страница
    story.append(Paragraph(presentation.title, title_style))
    story.append(Spacer(1, 20))
    story.append(Paragraph("Создано в SayDeck", styles['Normal']))
    story.append(Spacer(1, 40))
    
    # Добавляем слайды
    for i, slide_data in enumerate(presentation.content.get("slides", [])):
        if i > 0:  # Новая страница для каждого слайда кроме первого
            story.append(PageBreak())
        
        # Заголовок слайда
        story.append(Paragraph(f"Слайд {i+1}: {slide_data.get('title', '')}", slide_title_style))
        
        # Содержимое слайда
        content = slide_data.get('content', '').replace('\n', '<br/>')
        story.append(Paragraph(content, content_style))
        
        story.append(Spacer(1, 20))
    
    # Генерируем PDF
    doc.build(story)
    temp_file.close()
    
    # Возвращаем файл
    safe_title = "".join(c for c in presentation.title if c.isalnum() or c in (' ', '-', '_')).rstrip()
    filename = f"{safe_title}.pdf"
    
    return FileResponse(
        temp_file.name,
        media_type='application/pdf',
        filename=filename,
        background=lambda: os.unlink(temp_file.name)
    )

@router.get("/presentations/{presentation_id}/html")
async def export_to_html(
    presentation_id: int,
    current_user: User = Depends(get_current_user),
    session: AsyncSession = Depends(get_session)
):
    """Экспорт презентации в HTML формат"""
    
    # Получаем презентацию
    result = await session.execute(
        select(Presentation).where(
            Presentation.id == presentation_id,
            Presentation.user_id == current_user.id
        )
    )
    presentation = result.scalars().first()
    
    if not presentation:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Презентация не найдена"
        )
    
    # Генерируем HTML слайды
    slides_html = ""
    for i, slide in enumerate(presentation.content.get("slides", [])):
        slides_html += f"""
        <section class="slide" id="slide-{i}">
            <h2>{slide.get('title', '')}</h2>
            <div class="slide-content">
                <p>{slide.get('content', '').replace(chr(10), '</p><p>')}</p>
            </div>
        </section>
        """
    
    # HTML шаблон с reveal.js стилем
    html_content = f"""
    <!DOCTYPE html>
    <html lang="ru">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>{presentation.title}</title>
        <style>
            body {{
                font-family: 'Arial', sans-serif;
                margin: 0;
                padding: 0;
                background: linear-gradient(45deg, #1e3c72, #2a5298);
                color: white;
                overflow: hidden;
            }}
            
            .presentation {{
                width: 100vw;
                height: 100vh;
                display: flex;
                align-items: center;
                justify-content: center;
            }}
            
            .slide {{
                display: none;
                text-align: center;
                max-width: 900px;
                padding: 60px;
                background: rgba(255, 255, 255, 0.1);
                backdrop-filter: blur(10px);
                border-radius: 20px;
                box-shadow: 0 8px 32px rgba(31, 38, 135, 0.37);
                border: 1px solid rgba(255, 255, 255, 0.18);
            }}
            
            .slide.active {{
                display: block;
                animation: fadeIn 0.5s ease-in-out;
            }}
            
            @keyframes fadeIn {{
                from {{ opacity: 0; transform: translateY(20px); }}
                to {{ opacity: 1; transform: translateY(0); }}
            }}
            
            .slide h2 {{
                font-size: 3em;
                margin-bottom: 30px;
                color: #fff;
                text-shadow: 2px 2px 4px rgba(0,0,0,0.5);
            }}
            
            .slide-content {{
                font-size: 1.5em;
                line-height: 1.6;
                color: rgba(255, 255, 255, 0.9);
            }}
            
            .navigation {{
                position: fixed;
                bottom: 30px;
                left: 50%;
                transform: translateX(-50%);
                z-index: 1000;
            }}
            
            .nav-button {{
                background: rgba(255, 255, 255, 0.2);
                border: none;
                color: white;
                padding: 15px 25px;
                margin: 0 10px;
                border-radius: 30px;
                cursor: pointer;
                font-size: 16px;
                transition: all 0.3s;
                backdrop-filter: blur(10px);
            }}
            
            .nav-button:hover {{
                background: rgba(255, 255, 255, 0.3);
                transform: translateY(-2px);
            }}
            
            .nav-button:disabled {{
                opacity: 0.5;
                cursor: not-allowed;
            }}
            
            .slide-counter {{
                position: fixed;
                top: 30px;
                right: 30px;
                background: rgba(0, 0, 0, 0.5);
                padding: 10px 20px;
                border-radius: 20px;
                font-size: 14px;
            }}
            
            .progress-bar {{
                position: fixed;
                bottom: 0;
                left: 0;
                height: 4px;
                background: linear-gradient(90deg, #00c6ff, #0072ff);
                transition: width 0.3s ease;
                z-index: 1000;
            }}
        </style>
    </head>
    <body>
        <div class="presentation">
            {slides_html}
        </div>
        
        <div class="slide-counter">
            <span id="current-slide">1</span> / <span id="total-slides">{len(presentation.content.get('slides', []))}</span>
        </div>
        
        <div class="progress-bar" id="progress-bar"></div>
        
        <div class="navigation">
            <button class="nav-button" id="prev-btn" onclick="previousSlide()">← Назад</button>
            <button class="nav-button" id="next-btn" onclick="nextSlide()">Вперед →</button>
        </div>

        <script>
            let currentSlide = 0;
            const slides = document.querySelectorAll('.slide');
            const totalSlides = slides.length;
            
            function showSlide(n) {{
                slides.forEach(slide => slide.classList.remove('active'));
                
                if (n >= totalSlides) currentSlide = totalSlides - 1;
                if (n < 0) currentSlide = 0;
                
                slides[currentSlide].classList.add('active');
                
                // Обновляем счетчик
                document.getElementById('current-slide').textContent = currentSlide + 1;
                
                // Обновляем прогресс-бар
                const progress = ((currentSlide + 1) / totalSlides) * 100;
                document.getElementById('progress-bar').style.width = progress + '%';
                
                // Обновляем кнопки
                document.getElementById('prev-btn').disabled = currentSlide === 0;
                document.getElementById('next-btn').disabled = currentSlide === totalSlides - 1;
            }}
            
            function nextSlide() {{
                if (currentSlide < totalSlides - 1) {{
                    currentSlide++;
                    showSlide(currentSlide);
                }}
            }}
            
            function previousSlide() {{
                if (currentSlide > 0) {{
                    currentSlide--;
                    showSlide(currentSlide);
                }}
            }}
            
            // Управление клавиатурой
            document.addEventListener('keydown', function(e) {{
                switch(e.key) {{
                    case 'ArrowRight':
                    case ' ':
                        e.preventDefault();
                        nextSlide();
                        break;
                    case 'ArrowLeft':
                        e.preventDefault();
                        previousSlide();
                        break;
                    case 'Home':
                        e.preventDefault();
                        currentSlide = 0;
                        showSlide(currentSlide);
                        break;
                    case 'End':
                        e.preventDefault();
                        currentSlide = totalSlides - 1;
                        showSlide(currentSlide);
                        break;
                }}
            }});
            
            // Инициализация
            showSlide(0);
        </script>
    </body>
    </html>
    """
    
    # Сохраняем HTML во временный файл
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".html", mode='w', encoding='utf-8')
    temp_file.write(html_content)
    temp_file.close()
    
    # Возвращаем файл
    safe_title = "".join(c for c in presentation.title if c.isalnum() or c in (' ', '-', '_')).rstrip()
    filename = f"{safe_title}.html"
    
    return FileResponse(
        temp_file.name,
        media_type='text/html',
        filename=filename,
        background=lambda: os.unlink(temp_file.name)
    )
