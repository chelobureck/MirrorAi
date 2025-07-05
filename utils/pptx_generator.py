"""
PPTX Generator - создание презентаций в формате PowerPoint
"""
from pptx import Presentation as PPTXPresentation
import tempfile
import re
from typing import Dict, Any

async def generate_presentation_pptx(content: Dict[str, Any]) -> str:
    """Создает PPTX файл из структуры презентации"""
    prs = PPTXPresentation()
    
    for slide_data in content.get("slides", []):
        slide_layout = prs.slide_layouts[1]  # обычный слайд
        pptx_slide = prs.slides.add_slide(slide_layout)
        
        # Очищаем HTML теги из заголовка и контента для PPTX
        title = clean_html_tags(slide_data.get("title", ""))
        content_text = clean_html_tags(slide_data.get("content", ""))
        
        pptx_slide.shapes.title.text = title
        if len(pptx_slide.placeholders) > 1:
            pptx_slide.placeholders[1].text = content_text
    
    # Сохраняем во временный файл
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name 

def clean_html_tags(text: str) -> str:
    """Удаляет HTML теги из текста для использования в PPTX"""
    if not text:
        return ""
    
    # Заменяем основные HTML теги на текстовые эквиваленты
    text = text.replace('<br>', '\n')
    text = text.replace('<br/>', '\n')
    text = text.replace('<br />', '\n')
    
    # Удаляем все остальные HTML теги
    clean_text = re.sub(r'<[^>]+>', '', text)
    
    # Убираем лишние пробелы и переносы
    clean_text = re.sub(r'\s+', ' ', clean_text).strip()
    
    return clean_text
