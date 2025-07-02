import openai
from config.settings import get_settings
from typing import Dict, Any
from pptx import Presentation as PPTXPresentation
import tempfile
import json

settings = get_settings()

def get_openai_client():
    """Ленивая инициализация OpenAI клиента"""
    if not hasattr(get_openai_client, '_client'):
        get_openai_client._client = openai.OpenAI(api_key=settings.OPENAI_API_KEY)
        print(f"✓ OpenAI client initialized with key: {settings.OPENAI_API_KEY[:10]}...")
    return get_openai_client._client

async def transcribe_audio(audio_file_path: str) -> str:
    """Транскрибирует аудио файл с помощью Whisper API"""
    client = get_openai_client()
    with open(audio_file_path, "rb") as audio_file:
        transcript = client.audio.transcriptions.create(
            model="whisper-1",
            file=audio_file
        )
    return transcript.text

async def generate_presentation_structure(
    text: str, 
    language: str = "ru", 
    animation: bool = False, 
    slides_count: int = 5
) -> Dict[str, Any]:
    """Генерирует структуру презентации с помощью GPT-4"""
    prompt = f"""
    Создай структуру презентации на основе следующего текста на языке: {language}. 
    Количество слайдов: {slides_count}
    Анимация: {"включена" if animation else "отключена"}
    
    Верни ТОЛЬКО валидный JSON с полями:
    - title: заголовок презентации
    - slides: массив слайдов, где каждый слайд имеет поля:
        - title: заголовок слайда
        - content: основной текст слайда
        - type: тип слайда (title, content, image, etc.)
    
    Текст: {text}
    """
    
    client = get_openai_client()
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "Ты - эксперт по созданию презентаций. Отвечай только валидным JSON."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7
    )
    
    try:
        return json.loads(response.choices[0].message.content)
    except json.JSONDecodeError:
        # Fallback structure if JSON parsing fails
        return {
            "title": "Сгенерированная презентация",
            "slides": [
                {
                    "title": f"Слайд {i+1}",
                    "content": f"Содержимое слайда {i+1}",
                    "type": "content"
                } for i in range(slides_count)
            ]
        } 

async def generate_presentation_pptx(content: dict) -> str:
    prs = PPTXPresentation()
    for slide in content.get("slides", []):
        slide_layout = prs.slide_layouts[1]  # обычный слайд
        pptx_slide = prs.slides.add_slide(slide_layout)
        pptx_slide.shapes.title.text = slide.get("title", "")
        pptx_slide.placeholders[1].text = slide.get("content", "")
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name 